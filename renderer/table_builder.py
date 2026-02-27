import re
from pptx.util import Cm, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.oxml.ns import qn
import lxml.etree as etree


COLOR_BLACK = RGBColor(0x00, 0x00, 0x00)
COLOR_WHITE = RGBColor(0xFF, 0xFF, 0xFF)
TABLE_ROW_HEIGHT_CM = 1.25
TABLE_EXTRA_ROW_CM = 0.95
TABLE_CELL_PADDING_CM = 0.203
TABLE_BORDER_PT = 1.0
TABLE_MAX_WIDTH_CM = 30.480
TABLE_TEXT_SPACING_CM = 0.847

_SEPARATOR_RE = re.compile(
    r"^\s*\|?\s*:?-{3,}:?\s*(\|\s*:?-{3,}:?\s*)+\|?\s*$"
)


def OxmlElement(tag):
    return etree.SubElement(etree.Element("dummy"), qn(tag))


def parse_markdown_table(table_text: str) -> tuple[list[str], list[list[str]]]:
    """
    마크다운 표 텍스트를 파싱한다.
    Returns: (headers, rows)
      headers: ["구분", "내용", "비고"]
      rows: [["1단계", "평가", "-"], ["2단계", "처치", "-"]]
    """
    lines = [line.strip() for line in table_text.strip().splitlines() if line.strip()]
    if not lines:
        return [], []

    def split_row(line: str) -> list[str]:
        row = line.strip()
        if row.startswith("|"):
            row = row[1:]
        if row.endswith("|"):
            row = row[:-1]
        return [cell.strip() for cell in row.split("|")]

    headers = split_row(lines[0])
    col_count = len(headers)
    rows: list[list[str]] = []

    for idx, line in enumerate(lines[1:], start=1):
        if idx == 1 and _SEPARATOR_RE.match(line):
            continue
        cells = split_row(line)
        if len(cells) < col_count:
            cells.extend([""] * (col_count - len(cells)))
        elif len(cells) > col_count:
            cells = cells[:col_count]
        rows.append(cells)

    return headers, rows


def _count_lines(text: str) -> int:
    if not text:
        return 1
    return max(1, len(text.splitlines()))


def _row_height_cm(max_line_count: int) -> float:
    lines = max(1, max_line_count)
    return TABLE_ROW_HEIGHT_CM + (lines - 1) * TABLE_EXTRA_ROW_CM


def _calculate_row_heights(headers: list[str], rows: list[list[str]]) -> list[float]:
    header_lines = max((_count_lines(cell) for cell in headers), default=1)
    heights = [_row_height_cm(header_lines)]
    for row in rows:
        max_lines = max((_count_lines(cell) for cell in row), default=1)
        heights.append(_row_height_cm(max_lines))
    return heights


def _set_cell_border(cell) -> None:
    tc = cell._tc
    tcPr = tc.get_or_add_tcPr()
    for edge in ["a:lnL", "a:lnR", "a:lnT", "a:lnB"]:
        ln = tcPr.find(qn(edge))
        if ln is None:
            ln = OxmlElement(edge)
            tcPr.append(ln)
        else:
            for child in list(ln):
                ln.remove(child)
        ln.set("w", str(int(Pt(TABLE_BORDER_PT))))
        ln.set("cap", "flat")
        ln.set("cmpd", "sng")
        solid_fill = OxmlElement("a:solidFill")
        srgb_clr = OxmlElement("a:srgbClr")
        srgb_clr.set("val", "000000")
        solid_fill.append(srgb_clr)
        ln.append(solid_fill)


def _set_cell_text(
    cell,
    text: str,
    font_name: str,
    font_size: int,
    bold: bool,
) -> None:
    text_frame = cell.text_frame
    text_frame.clear()
    text_frame.word_wrap = True

    paragraph = text_frame.paragraphs[0]
    paragraph.text = text

    for p in text_frame.paragraphs:
        p.alignment = PP_ALIGN.CENTER
        if not p.runs:
            p.add_run()
        for run in p.runs:
            run.font.name = font_name
            run.font.size = Pt(font_size)
            run.font.bold = bold
            run.font.color.rgb = COLOR_BLACK


def add_table_to_slide(
    slide,
    table_text: str,
    left_cm: float,
    top_cm: float,
    max_width_cm: float = 30.480,
    font_name: str = "Malgun Gothic",
    font_size: int = 14,
) -> float:
    """
    슬라이드에 표를 추가한다.
    Returns: 표의 총 높이 (cm).
    """
    headers, rows = parse_markdown_table(table_text)
    if not headers:
        return 0.0

    cols_count = len(headers)
    rows_count = len(rows) + 1
    max_width_cm = min(max_width_cm, TABLE_MAX_WIDTH_CM)
    col_width_cm = max_width_cm / cols_count

    row_heights = _calculate_row_heights(headers, rows)
    total_height = sum(row_heights)

    table_shape = slide.shapes.add_table(
        rows_count,
        cols_count,
        Cm(left_cm),
        Cm(top_cm),
        Cm(max_width_cm),
        Cm(total_height),
    )
    table = table_shape.table

    col_width_emu = Emu(int(Cm(col_width_cm)))
    for col_idx in range(cols_count):
        table.columns[col_idx].width = col_width_emu

    for row_idx, row_height_cm in enumerate(row_heights):
        table.rows[row_idx].height = Emu(int(Cm(row_height_cm)))

    for row_idx in range(rows_count):
        for col_idx in range(cols_count):
            cell = table.cell(row_idx, col_idx)
            cell.fill.solid()
            cell.fill.fore_color.rgb = COLOR_WHITE
            _set_cell_border(cell)

            cell.vertical_anchor = MSO_ANCHOR.MIDDLE
            padding = Cm(TABLE_CELL_PADDING_CM)
            cell.margin_left = padding
            cell.margin_right = padding
            cell.margin_top = padding
            cell.margin_bottom = padding

            if row_idx == 0:
                text = headers[col_idx]
                bold = True
            else:
                text = rows[row_idx - 1][col_idx]
                bold = False
            _set_cell_text(cell, text, font_name, font_size, bold)

    try:
        table_shape.shadow.inherit = False
    except (NotImplementedError, AttributeError, Exception):
        pass

    return total_height


def estimate_table_height(table_text: str) -> float:
    """표의 예상 높이(cm)를 계산한다. 표 전후 간격 포함."""
    headers, rows = parse_markdown_table(table_text)
    if not headers:
        return 0.0
    total_height = sum(_calculate_row_heights(headers, rows))
    return total_height + (TABLE_TEXT_SPACING_CM * 2)


if __name__ == "__main__":
    from pathlib import Path
    from pptx import Presentation

    test_table = """| 구분 | 내용 | 비고 |
|------|------|------|
| 1단계 | 평가 | 초기 |
| 2단계 | 처치 | 현장 |
| 3단계 | 이송 | 병원 |"""

    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    add_table_to_slide(slide, test_table, left_cm=1.0, top_cm=1.0)

    output_path = Path("/Users/chungji/Documents/md2pptx/test_table.pptx")
    output_path.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(output_path))
