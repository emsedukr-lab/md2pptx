import sys
import re

sys.path.insert(0, '/Users/chungji/Documents/md2pptx')

from pptx import Presentation
from pptx.util import Cm, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.oxml.ns import qn

from config import (
    SlideData,
    UserSettings,
    SLIDE_WIDTH_CM,
    SLIDE_HEIGHT_CM,
    TITLE_LEFT,
    TITLE_TOP,
    TITLE_WIDTH,
    TITLE_HEIGHT,
    PAGE_NUM_LEFT,
    PAGE_NUM_TOP,
    PAGE_NUM_WIDTH,
    PAGE_NUM_HEIGHT,
    SEPARATOR_LEFT,
    SEPARATOR_TOP,
    SEPARATOR_WIDTH,
    SEPARATOR_WEIGHT_PT,
    BODY_LEFT,
    BODY_TOP,
    BODY_WIDTH,
    BODY_HEIGHT,
    COLOR_BLACK,
    DEFAULT_FONT,
    CODE_FONT,
    DEFAULT_TITLE_SIZE,
    DEFAULT_BODY_SIZE,
    DEFAULT_TABLE_SIZE,
    DEFAULT_CODE_BLOCK_SIZE,
    DEFAULT_PAGE_NUM_SIZE,
)
from renderer.text_formatter import format_body_text
from renderer.table_builder import add_table_to_slide, estimate_table_height


_TABLE_SEPARATOR_RE = re.compile(r"^\s*\|?\s*:?-{3,}:?\s*(\|\s*:?-{3,}:?\s*)+\|?\s*$")
_TEXT_BLOCK_MIN_HEIGHT_CM = 1.0
_TEXT_TABLE_GAP_CM = 0.35
_PT_TO_CM = 0.0352778


def _rgb_from_hex(hex_color: str) -> RGBColor:
    color = hex_color.strip().lstrip("#")
    if len(color) != 6:
        return RGBColor(0x00, 0x00, 0x00)
    return RGBColor(int(color[0:2], 16), int(color[2:4], 16), int(color[4:6], 16))


def _set_run_font(run, font_name: str, font_size_pt: int, bold: bool = False) -> None:
    run.font.name = font_name
    run.font.size = Pt(font_size_pt)
    run.font.bold = bold
    run.font.color.rgb = _rgb_from_hex(COLOR_BLACK)

    r_pr = run._r.get_or_add_rPr()
    ea = r_pr.find(qn("a:ea"))
    if ea is None:
        import lxml.etree as etree
        ea = etree.SubElement(r_pr, qn("a:ea"))
    ea.set("typeface", font_name)


def _first_run(paragraph):
    if paragraph.runs:
        return paragraph.runs[0]
    return paragraph.add_run()


def _is_pipe_row(line: str) -> bool:
    stripped = line.strip()
    return bool(stripped) and stripped.startswith("|") and stripped.endswith("|")


def _is_markdown_table_block(lines: list[str]) -> bool:
    if len(lines) < 3:
        return False
    if not all(_is_pipe_row(line) for line in lines):
        return False
    return bool(_TABLE_SEPARATOR_RE.match(lines[1]))


def _estimate_text_height_cm(text: str, body_size_pt: int) -> float:
    lines = text.splitlines() if text else [""]
    line_count = max(1, len(lines))
    line_height_cm = (body_size_pt * 1.5) * _PT_TO_CM
    paragraph_gap_cm = (6 * _PT_TO_CM) * 0.4
    return max(_TEXT_BLOCK_MIN_HEIGHT_CM, (line_count * line_height_cm) + paragraph_gap_cm)


def _append_body_text_box(
    slide,
    text: str,
    top_cm: float,
    height_cm: float,
    settings: UserSettings,
) -> None:
    body_box = slide.shapes.add_textbox(
        Cm(BODY_LEFT),
        Cm(top_cm),
        Cm(BODY_WIDTH),
        Cm(height_cm),
    )
    body_box.text_frame.auto_size = None
    format_body_text(
        body_box.text_frame,
        text,
        font_name=settings.font_name,
        body_size=settings.body_size,
        code_font=settings.code_font,
        code_block_size=DEFAULT_CODE_BLOCK_SIZE,
        line_spacing_multiplier=settings.line_spacing_multiplier,
    )


def create_presentation(
    slides_data: list[SlideData],
    settings: UserSettings = None,
) -> Presentation:
    prs = Presentation()
    prs.slide_width = Cm(SLIDE_WIDTH_CM)
    prs.slide_height = Cm(SLIDE_HEIGHT_CM)

    settings = settings or UserSettings(
        font_name=DEFAULT_FONT,
        title_size=DEFAULT_TITLE_SIZE,
        body_size=DEFAULT_BODY_SIZE,
        table_size=DEFAULT_TABLE_SIZE,
        code_font_separate=False,
    )

    for slide_data in slides_data:
        for split_slide_data in _split_if_overflow(slide_data):
            _build_slide(prs, split_slide_data, settings)

    return prs


def _build_slide(
    prs: Presentation,
    slide_data: SlideData,
    settings: UserSettings,
) -> None:
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)

    title_box = slide.shapes.add_textbox(
        Cm(TITLE_LEFT),
        Cm(TITLE_TOP),
        Cm(TITLE_WIDTH),
        Cm(TITLE_HEIGHT),
    )
    tf = title_box.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.BOTTOM
    p = tf.paragraphs[0]
    p.text = slide_data.title
    p.alignment = PP_ALIGN.LEFT
    run = _first_run(p)
    _set_run_font(run, settings.font_name, settings.title_size, bold=True)

    if slide_data.page_number:
        page_box = slide.shapes.add_textbox(
            Cm(PAGE_NUM_LEFT),
            Cm(PAGE_NUM_TOP),
            Cm(PAGE_NUM_WIDTH),
            Cm(PAGE_NUM_HEIGHT),
        )
        tf = page_box.text_frame
        tf.word_wrap = False
        tf.vertical_anchor = MSO_ANCHOR.BOTTOM
        p = tf.paragraphs[0]
        p.text = slide_data.page_number
        p.alignment = PP_ALIGN.RIGHT
        run = _first_run(p)
        _set_run_font(run, settings.font_name, DEFAULT_PAGE_NUM_SIZE, bold=False)

    from pptx.enum.shapes import MSO_CONNECTOR

    left = Cm(SEPARATOR_LEFT)
    top = Cm(SEPARATOR_TOP)
    width = Cm(SEPARATOR_WIDTH)
    connector = slide.shapes.add_connector(
        MSO_CONNECTOR.STRAIGHT,
        left,
        top,
        Emu(int(left + width)),
        top,
    )
    connector.line.color.rgb = _rgb_from_hex(COLOR_BLACK)
    connector.line.width = Pt(SEPARATOR_WEIGHT_PT)

    if hasattr(connector, "shadow"):
        try:
            connector.shadow.inherit = False
        except Exception:
            pass

    body_text = slide_data.body_text or ""
    if not body_text.strip():
        pass
    else:
        segments = _split_body_segments(body_text)
        has_table = any(seg["type"] == "table" for seg in segments)

        if not has_table:
            body_box = slide.shapes.add_textbox(
                Cm(BODY_LEFT),
                Cm(BODY_TOP),
                Cm(BODY_WIDTH),
                Cm(BODY_HEIGHT),
            )
            body_box.text_frame.auto_size = None
            format_body_text(
                body_box.text_frame,
                body_text,
                font_name=settings.font_name,
                body_size=settings.body_size,
                code_font=settings.code_font,
                code_block_size=DEFAULT_CODE_BLOCK_SIZE,
                line_spacing_multiplier=settings.line_spacing_multiplier,
            )
        else:
            current_top = BODY_TOP
            body_bottom = BODY_TOP + BODY_HEIGHT

            for segment in segments:
                if current_top >= body_bottom:
                    break

                segment_type = segment.get("type")
                segment_content = segment.get("content", "")

                if segment_type == "text":
                    if not segment_content.strip():
                        continue

                    estimated_height = _estimate_text_height_cm(segment_content, settings.body_size)
                    available_height = max(0.0, body_bottom - current_top)
                    box_height = min(estimated_height, available_height)
                    if box_height <= 0:
                        break

                    _append_body_text_box(slide, segment_content, current_top, box_height, settings)
                    current_top += box_height + _TEXT_TABLE_GAP_CM
                    continue

                if segment_type == "table":
                    if not segment_content.strip():
                        continue

                    estimated_table_height = estimate_table_height(segment_content)
                    available_height = max(0.0, body_bottom - current_top)
                    if available_height <= 0:
                        break

                    if estimated_table_height > available_height:
                        break

                    actual_table_height = add_table_to_slide(
                        slide,
                        segment_content,
                        left_cm=BODY_LEFT,
                        top_cm=current_top,
                        max_width_cm=BODY_WIDTH,
                        font_name=settings.font_name,
                        font_size=settings.table_size,
                    )
                    if actual_table_height <= 0:
                        continue

                    consumed_height = max(actual_table_height, estimated_table_height)
                    current_top += consumed_height

    if slide_data.speaker_notes:
        notes_slide = slide.notes_slide
        notes_tf = notes_slide.notes_text_frame
        notes_tf.text = slide_data.speaker_notes


def _split_body_segments(body_text: str) -> list[dict]:
    """
    body_text를 텍스트/표 세그먼트로 분리.
    Returns: [{"type": "text", "content": "..."}, {"type": "table", "content": "..."}]
    """
    if not body_text:
        return []

    lines = body_text.splitlines()
    segments: list[dict] = []
    text_buffer: list[str] = []
    i = 0

    while i < len(lines):
        line = lines[i]

        if _is_pipe_row(line):
            candidate: list[str] = []
            j = i
            while j < len(lines) and _is_pipe_row(lines[j]):
                candidate.append(lines[j])
                j += 1

            if _is_markdown_table_block(candidate):
                text_content = "\n".join(text_buffer).strip("\n")
                if text_content.strip():
                    segments.append({"type": "text", "content": text_content})
                segments.append({"type": "table", "content": "\n".join(candidate)})
                text_buffer = []
                i = j
                continue

            text_buffer.extend(candidate)
            i = j
            continue

        text_buffer.append(line)
        i += 1

    remaining_text = "\n".join(text_buffer).strip("\n")
    if remaining_text.strip():
        segments.append({"type": "text", "content": remaining_text})

    if not segments and body_text.strip():
        segments.append({"type": "text", "content": body_text})

    return segments


def _split_if_overflow(slide_data: SlideData, max_lines: int = 25) -> list[SlideData]:
    body_text = slide_data.body_text or ""
    lines = body_text.splitlines()

    if len(lines) <= max_lines:
        return [slide_data]

    split_slides: list[SlideData] = []
    for idx in range(0, len(lines), max_lines):
        chunk_lines = lines[idx : idx + max_lines]
        chunk_index = idx // max_lines

        chunk_title = slide_data.title if chunk_index == 0 else f"{slide_data.title} (cont.)"
        chunk_page_number = slide_data.page_number if chunk_index == 0 else ""
        chunk_notes = slide_data.speaker_notes if chunk_index == 0 else ""

        split_slides.append(
            SlideData(
                title=chunk_title,
                page_number=chunk_page_number,
                body_text="\n".join(chunk_lines).strip("\n"),
                speaker_notes=chunk_notes,
            )
        )

    return split_slides


if __name__ == "__main__":
    from parser.md_parser import parse_markdown

    test_markdown = """# [1]: 학습목표 (교재 p.3)

**[핵심 메시지]**
본 장은 병원 전 외상 처치의 역사를 다룬다. 외상 처치의 세 단계를 학습한다.

[본문]
- 병원 전 외상 처치의 역사와 발전을 이해할 수 있다.
  - 고대부터 현대까지의 발전 과정
    - 나폴레옹 시대의 군의관 체계
- 외상 처치의 세 가지 단계를 이해할 수 있다.
- **외상성 손상**의 인적 및 재정적 영향의 중대성을 인식할 수 있다.

# [5]: 외상의 분류 (교재 p.102~103)

**[핵심 메시지]**
외상은 둔상과 관통상으로 분류된다.

[본문]
- 둔상(blunt trauma): 표면 손상 없이 내부 장기 손상
- 관통상(penetrating trauma): 피부를 뚫고 들어가는 손상

| 구분 | 설명 | 예시 |
|------|------|------|
| 둔상 | 내부 손상 | 교통사고 |
| 관통상 | 외부→내부 | 총상 |
"""

    parsed_slides = parse_markdown(test_markdown)
    slides_data = [
        SlideData(
            title=slide.title,
            page_number=slide.page_number,
            body_text=slide.body_text,
            speaker_notes=slide.speaker_notes,
        )
        for slide in parsed_slides
    ]

    presentation = create_presentation(slides_data)
    presentation.save("/Users/chungji/Documents/md2pptx/test_output.pptx")
