import re
import math
from pptx.util import Cm, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR


COLOR_BLACK = RGBColor(0x00, 0x00, 0x00)
COLOR_GRAY_100 = RGBColor(0xF3, 0xF4, 0xF6)
BULLET_INDENT_CM = 1.2
DEFAULT_LINE_SPACING_MULTIPLIER = 1.5
PARAGRAPH_AFTER_PT = 6
_BASE_MAX_CHARS = 65
_BASE_BREAK_AT = 58
_LEVEL_CHAR_REDUCTION = 10

# P1: 수평 구분선 패턴 (---, ***, ___)
_HORIZONTAL_RULE_RE = re.compile(r"^\s*[-]{3,}\s*$|^\s*[*]{3,}\s*$|^\s*[_]{3,}\s*$")


INLINE_PATTERNS = [
    ("inline_code", re.compile(r"`([^`]+)`")),
    ("bold_italic", re.compile(r"\*\*\*([\s\S]+?)\*\*\*")),
    ("bold", re.compile(r"\*\*([\s\S]+?)\*\*")),
    ("italic", re.compile(r"\*([\s\S]+?)\*")),
    ("strikethrough", re.compile(r"~~([\s\S]+?)~~")),
    ("latex_inline", re.compile(r"\$([^$]+)\$")),
]


def _create_paragraph(text_frame, use_first_paragraph: bool):
    paragraph = text_frame.paragraphs[0] if use_first_paragraph else text_frame.add_paragraph()
    paragraph.alignment = PP_ALIGN.LEFT
    return paragraph


def _set_run_background(run, color: RGBColor) -> None:
    from pptx.oxml.xmlchemy import OxmlElement

    r_pr = run._r.get_or_add_rPr()
    for child in list(r_pr):
        if child.tag.endswith("}highlight"):
            r_pr.remove(child)

    highlight = OxmlElement("a:highlight")
    srgb = OxmlElement("a:srgbClr")
    srgb.set("val", f"{color[0]:02X}{color[1]:02X}{color[2]:02X}")
    highlight.append(srgb)
    r_pr.append(highlight)


def _configure_run(run, font_name: str, font_size: int, *, bold=False, italic=False, strike=False) -> None:
    run.font.name = font_name
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.strike = strike
    run.font.color.rgb = COLOR_BLACK


def _tokenize_inline(text: str):
    tokens = []
    cursor = 0

    while cursor < len(text):
        best_match = None
        best_kind = None
        best_priority = None

        for priority, (kind, pattern) in enumerate(INLINE_PATTERNS):
            match = pattern.search(text, cursor)
            if match is None:
                continue

            if best_match is None:
                best_match = match
                best_kind = kind
                best_priority = priority
                continue

            if match.start() < best_match.start():
                best_match = match
                best_kind = kind
                best_priority = priority
                continue

            if match.start() == best_match.start() and priority < best_priority:
                best_match = match
                best_kind = kind
                best_priority = priority

        if best_match is None:
            tokens.append(("text", text[cursor:]))
            break

        if best_match.start() > cursor:
            tokens.append(("text", text[cursor : best_match.start()]))

        tokens.append((best_kind, best_match.group(1)))
        cursor = best_match.end()

    return tokens


def _set_list_indentation(paragraph, level: int) -> None:
    p_pr = paragraph._p.get_or_add_pPr()
    p_pr.indent = Emu(0)
    p_pr.marL = Cm(BULLET_INDENT_CM * level)


def _append_code_block(paragraph, code_lines, code_font: str, code_block_size: int, line_spacing_multiplier: float = 1.5) -> None:
    paragraph.space_after = Pt(PARAGRAPH_AFTER_PT)
    paragraph.line_spacing = Pt(int(code_block_size * line_spacing_multiplier))

    if not code_lines:
        run = paragraph.add_run()
        run.text = "  "
        _configure_run(run, code_font, code_block_size)
        _set_run_background(run, COLOR_GRAY_100)
        return

    for idx, line in enumerate(code_lines):
        run = paragraph.add_run()
        run.text = f"  {line}  "
        _configure_run(run, code_font, code_block_size)
        _set_run_background(run, COLOR_GRAY_100)
        if idx < len(code_lines) - 1:
            _add_soft_break(paragraph)


def _add_soft_break(paragraph) -> None:
    """paragraph 내에 soft line break (Shift+Enter)를 추가한다."""
    from pptx.oxml.ns import qn
    import lxml.etree as etree
    etree.SubElement(paragraph._p, qn("a:br"))


def _apply_soft_line_break(paragraph, text: str, font_name: str, font_size: int, code_font: str, max_chars: int = _BASE_MAX_CHARS, break_at: int = _BASE_BREAK_AT) -> None:
    """P6: 레벨별 차등 줄바꿈. max_chars 이상 줄을 break_at 어절 기준으로 soft break 후 2칸 들여쓰기."""
    if len(text) < max_chars:
        apply_inline_formatting(paragraph, text, font_name, font_size, code_font)
        return

    words = text.split()
    if not words:
        apply_inline_formatting(paragraph, text, font_name, font_size, code_font)
        return

    cumulative = 0
    split_idx = len(words)
    for idx, word in enumerate(words):
        added = len(word) if idx == 0 else len(word) + 1
        cumulative += added
        if cumulative > break_at:
            split_idx = idx + 1
            break

    first_part = " ".join(words[:split_idx])
    rest_part = " ".join(words[split_idx:])

    apply_inline_formatting(paragraph, first_part, font_name, font_size, code_font)

    if rest_part:
        _add_soft_break(paragraph)
        # P4: 2칸 공백 들여쓰기 후 나머지 텍스트
        _apply_soft_line_break(paragraph, "  " + rest_part, font_name, font_size, code_font, max_chars, break_at)


def format_body_text(
    text_frame,
    markdown_text: str,
    font_name: str = "Malgun Gothic",
    body_size: int = 16,
    code_font: str = "Malgun Gothic",
    code_block_size: int = 14,
    line_spacing_multiplier: float = DEFAULT_LINE_SPACING_MULTIPLIER,
) -> None:
    text_frame.clear()
    text_frame.word_wrap = True
    text_frame.vertical_anchor = MSO_ANCHOR.TOP

    lines = markdown_text.splitlines()
    in_code_block = False
    code_lines = []
    use_first_paragraph = True

    for line in lines:
        fence_match = re.match(r"^\s*```", line)
        if fence_match:
            if in_code_block:
                paragraph = _create_paragraph(text_frame, use_first_paragraph)
                use_first_paragraph = False
                _append_code_block(paragraph, code_lines, code_font, code_block_size, line_spacing_multiplier)
                in_code_block = False
                code_lines = []
            else:
                in_code_block = True
                code_lines = []
            continue

        if in_code_block:
            code_lines.append(line)
            continue

        # P1: 수평 구분선 무시 (---, ***, ___)
        if _HORIZONTAL_RULE_RE.match(line):
            continue

        if line.strip() == "":
            paragraph = _create_paragraph(text_frame, use_first_paragraph)
            use_first_paragraph = False
            set_line_spacing(paragraph, body_size, line_spacing_multiplier)
            continue

        bullet_match = re.match(r"^(\s*)[-*]\s+(.*)$", line)
        if bullet_match:
            leading_spaces = len(bullet_match.group(1))
            level = math.ceil(leading_spaces / 2)
            item_text = bullet_match.group(2)

            paragraph = _create_paragraph(text_frame, use_first_paragraph)
            use_first_paragraph = False
            paragraph.level = level
            paragraph.space_after = Pt(PARAGRAPH_AFTER_PT)
            _set_list_indentation(paragraph, level)

            # P3: 화살표(→)로 시작하면 글머리 기호 생략
            if not item_text.strip().startswith("→"):
                bullet_run = paragraph.add_run()
                bullet_run.text = "• "
                _configure_run(bullet_run, font_name, body_size)

            # P6: 레벨별 차등 줄바꿈 기준
            lvl_max = max(20, _BASE_MAX_CHARS - level * _LEVEL_CHAR_REDUCTION)
            lvl_brk = max(13, _BASE_BREAK_AT - level * _LEVEL_CHAR_REDUCTION)
            _apply_soft_line_break(paragraph, item_text, font_name, body_size, code_font, lvl_max, lvl_brk)
            set_line_spacing(paragraph, body_size, line_spacing_multiplier)
            continue

        ordered_match = re.match(r"^(\s*)(\d+)\.\s+(.*)$", line)
        if ordered_match:
            leading_spaces = len(ordered_match.group(1))
            level = math.ceil(leading_spaces / 2)
            order_number = ordered_match.group(2)
            item_text = ordered_match.group(3)

            paragraph = _create_paragraph(text_frame, use_first_paragraph)
            use_first_paragraph = False
            paragraph.level = level
            paragraph.space_after = Pt(PARAGRAPH_AFTER_PT)
            _set_list_indentation(paragraph, level)

            order_run = paragraph.add_run()
            order_run.text = f"{order_number}. "
            _configure_run(order_run, font_name, body_size)

            # P6: 레벨별 차등 줄바꿈 기준
            lvl_max = max(20, _BASE_MAX_CHARS - level * _LEVEL_CHAR_REDUCTION)
            lvl_brk = max(13, _BASE_BREAK_AT - level * _LEVEL_CHAR_REDUCTION)
            _apply_soft_line_break(paragraph, item_text, font_name, body_size, code_font, lvl_max, lvl_brk)
            set_line_spacing(paragraph, body_size, line_spacing_multiplier)
            continue

        paragraph = _create_paragraph(text_frame, use_first_paragraph)
        use_first_paragraph = False
        _apply_soft_line_break(paragraph, line, font_name, body_size, code_font)
        set_line_spacing(paragraph, body_size, line_spacing_multiplier)

    if in_code_block:
        paragraph = _create_paragraph(text_frame, use_first_paragraph)
        _append_code_block(paragraph, code_lines, code_font, code_block_size, line_spacing_multiplier)


def apply_inline_formatting(
    paragraph,
    text: str,
    font_name: str,
    font_size: int,
    code_font: str,
) -> None:
    tokens = _tokenize_inline(text)

    for token_type, token_text in tokens:
        run = paragraph.add_run()
        run.text = token_text
        _configure_run(run, font_name, font_size)

        if token_type == "inline_code":
            run.font.name = code_font
            _set_run_background(run, COLOR_GRAY_100)
        elif token_type == "bold_italic":
            run.font.bold = True
            run.font.italic = True
        elif token_type == "bold":
            run.font.bold = True
        elif token_type == "italic":
            run.font.italic = True
        elif token_type == "strikethrough":
            run.font.strike = True
        elif token_type == "latex_inline":
            run.font.italic = True


def set_line_spacing(paragraph, font_size_pt: int, line_spacing_multiplier: float = DEFAULT_LINE_SPACING_MULTIPLIER) -> None:
    paragraph.line_spacing = Pt(int(font_size_pt * line_spacing_multiplier))
    paragraph.space_after = Pt(PARAGRAPH_AFTER_PT)


if __name__ == "__main__":
    from pptx import Presentation

    markdown_sample = """- 첫 번째 항목 **볼드** 포함
  - 두 번째 레벨 *이탤릭*
    - 세 번째 레벨 `코드`
- 일반 항목
1. 순서 리스트 첫 번째
2. 순서 리스트 두 번째

일반 텍스트 문단입니다. ***볼드이탤릭*** 테스트.
~~취소선~~ 테스트. $E=mc^2$ 수식 테스트.
"""

    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    textbox = slide.shapes.add_textbox(Cm(1.5), Cm(1.5), Cm(30), Cm(16))

    format_body_text(textbox.text_frame, markdown_sample)
    prs.save("/Users/chungji/Documents/md2pptx/test_formatter.pptx")
